# -*- coding: utf-8 -*-
"""Créer de nouveaux certificats laser au besoin."""

# Bibliothèques standard
import subprocess
import time

from pathlib import Path
from datetime import datetime as dt
from subprocess import run

# Bibliothèque PIPy
import schedule
import pptx
import keyring
import getpass

# Bibliothèques maison
from polygphys.outils.reseau.msforms import MSFormConfig, MSForm
from polygphys.outils.reseau import DisqueRéseau, OneDrive

# Définitions de classes

class SSTLaserCertificatsConfig(MSFormConfig):

    def default(self):
        return (Path(__file__).parent / 'nouveau_certificat.cfg').open().read()


class SSTLaserCertificatsForm(MSForm):

    def nettoyer(self, cadre):
        cadre = self.convertir_champs(cadre)
        cadre = cadre.astype({'matricule': int}, errors='ignore')

        courriels_manquants = cadre['courriel'] == 'anonymous'
        cadre.loc[courriels_manquants,
                  'courriel'] = cadre.loc[courriels_manquants, 'courriel2']
        cadre.courriel = cadre.courriel.fillna(
            cadre.courriel2).fillna('@polymtl.ca')
        cadre.nom = cadre.nom.fillna(cadre.nom2).fillna('anonyme')
        cadre.date = cadre.date.dt.date
        cadre.matricule = cadre.matricule.fillna(0)

        return cadre.loc[:, ['date', 'matricule', 'courriel', 'nom']]

    def action(self, cadre):
        for i, entrée in cadre.iterrows():
            chemin_cert = Path(__file__).parent / \
                self.config.get('certificats', 'chemin')
            cert = pptx.Presentation(chemin_cert)

            for forme in cert.slides[0].shapes:
                if forme.has_text_frame:
                    for par in forme.text_frame.paragraphs:
                        for ligne in par.runs:
                            if ligne.text == 'nom':
                                ligne.text = str(entrée.nom)
                            elif ligne.text == 'matricule':
                                ligne.text = str(entrée.matricule)
                            elif ligne.text.startswith('Date'):
                                date = dt.today()
                                ligne.text = f'Date: {date.year}-{date.month:02}'

            for disque in self.config.getlist('certificats', 'disques'):
                url = self.config.get(disque, 'url')
                chemin = self.config.getpath(disque, 'mount_point')
                drive = self.config.get(disque, 'drive')
                mode = self.config.get(disque, 'method')

                nom = self.config.get(disque, 'nom')
                mdp = keyring.get_password(
                    'system', f'polygphys.sst.laser.{disque}.{nom}')
                if mdp is None:
                    mdp = getpass.getpass('mdp: ')
                    keyring.set_password(
                        'system', f'polygphys.sst.laser.{disque}.{nom}', mdp)
                with DisqueRéseau(url, chemin, drive, nom, mdp, mode) as d:
                    sous_dossier = d / self.config.get(disque, 'chemin')
                    sous_dossier = d / self.config.get('certificats', 'ppt')
                    fichier = sous_dossier / f'{entrée.nom}.pptx'
                    cert.save(fichier)

                    fichier_pdf = fichier.parent.parent / 'pdf' / fichier.name
                    run(['unoconv',
                         '-f',
                         'pdf',
                         '-o',
                         str(fichier_pdf),
                         str(fichier)])

# Programme

chemin_config = Path('~').expanduser() / 'certificats_laser.cfg'
config = SSTLaserCertificatsConfig(chemin_config)

dossier = OneDrive('',
                   config.get('onedrive', 'organisation'),
                   config.get('onedrive', 'sous-dossier'),
                   partagé=True)
fichier = dossier / config.get('formulaire', 'nom')
config.set('formulaire', 'chemin', str(fichier))

formulaire = SSTLaserCertificatsForm(config)

exporteur = subprocess.Popen(['unoconv', '--listener'])

schedule.every().day.at('08:00').do(formulaire.mise_à_jour)

formulaire.mise_à_jour()
try:
    while True:
        schedule.run_pending()
        time.sleep(1)
except KeyboardInterrupt:
    print('On arrête.')
finally:
    exporteur.terminate()

print('Terminé.')
